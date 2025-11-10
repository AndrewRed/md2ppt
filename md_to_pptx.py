#!/usr/bin/env python3
"""
Улучшенный конвертер Markdown презентации в PowerPoint
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Цветовая схема
COLORS = {
    'primary': RGBColor(0, 51, 102),      # Темно-синий
    'accent': RGBColor(0, 102, 204),      # Синий
    'success': RGBColor(22, 163, 74),     # Зеленый
    'warning': RGBColor(217, 119, 6),     # Оранжевый
    'text': RGBColor(51, 51, 51),        # Темно-серый
    'light': RGBColor(102, 102, 102),     # Светло-серый
}

def parse_markdown_sections(md_content):
    """Парсит Markdown и извлекает разделы с подразделами"""
    sections = []
    current_section = {"title": "", "subsections": [], "content": []}
    
    lines = md_content.split('\n')
    current_subsection = None
    
    for i, line in enumerate(lines):
        line_stripped = line.strip()
        
        # Основной заголовок раздела (##)
        if line.startswith('##') and not line.startswith('###'):
            # Сохраняем предыдущий раздел
            if current_section["title"]:
                sections.append(current_section.copy())
            
            # Начинаем новый раздел
            title = line.replace('##', '').strip()
            title = re.sub(r'[0-9️⃣1️⃣2️⃣3️⃣4️⃣5️⃣6️⃣7️⃣8️⃣9️⃣🔟]', '', title).strip()
            title = re.sub(r'^\d+\.\s*', '', title)
            current_section = {"title": title, "subsections": [], "content": []}
            current_subsection = None
        
        # Подзаголовок (###)
        elif line.startswith('###'):
            if current_subsection:
                current_section["subsections"].append(current_subsection)
            current_subsection = {
                "title": line.replace('###', '').strip(),
                "content": []
            }
        
        # Обычный контент
        elif line_stripped and not line.startswith('---'):
            if current_subsection:
                current_subsection["content"].append(line)
            else:
                current_section["content"].append(line)
    
    # Сохраняем последний подраздел и раздел
    if current_subsection:
        current_section["subsections"].append(current_subsection)
    if current_section["title"]:
        sections.append(current_section)
    
    return sections

def clean_markdown_text(text, keep_emoji=True):
    """Очищает Markdown разметку из текста"""
    if not text:
        return ""
    
    # Убираем жирный текст (оставляем текст)
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Убираем курсив
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Убираем код блоки
    text = re.sub(r'```[\s\S]*?```', '', text)
    # Убираем инлайн код
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Убираем ссылки [текст](url)
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    
    return text.strip()

def extract_bullets(content_lines):
    """Извлекает маркированные списки из контента"""
    bullets = []
    skip_next = False
    
    for i, line in enumerate(content_lines):
        line = line.strip()
        if not line or line.startswith('---') or line.startswith('```'):
            skip_next = line.startswith('```')
            continue
        
        if skip_next and not line.startswith('```'):
            continue
        elif line.startswith('```'):
            skip_next = False
            continue
        
        # Маркированный список
        if re.match(r'^[-*+]\s+', line):
            bullet = re.sub(r'^[-*+]\s+', '', line)
            bullet = clean_markdown_text(bullet)
            if bullet:
                bullets.append(bullet)
        # Нумерованный список
        elif re.match(r'^\d+\.\s+', line):
            bullet = re.sub(r'^\d+\.\s+', '', line)
            bullet = clean_markdown_text(bullet)
            if bullet:
                bullets.append(bullet)
        # Заголовки подразделов (Сценарий, Преимущества и т.д.)
        elif re.match(r'^\*\*.*\*\*:', line):
            header = clean_markdown_text(line)
            bullets.append(header)
        # Обычный текст (если нет маркеров)
        elif line and not line.startswith('|'):  # Не таблица
            cleaned = clean_markdown_text(line)
            # Пропускаем очень короткие строки и примеры кода
            if cleaned and len(cleaned) > 15 and not cleaned.startswith('Пользователь:') and not cleaned.startswith('AI:'):
                bullets.append(cleaned)
    
    return bullets

def parse_table(content_lines):
    """Парсит Markdown таблицу"""
    table_data = []
    for line in content_lines:
        if '|' in line and not line.strip().startswith('|---'):
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if cells and not all(c == '-' for c in ''.join(cells)):
                table_data.append(cells)
    return table_data if len(table_data) > 1 else None

def create_slide_with_bullets(prs, title, bullets, max_bullets=7):
    """Создает слайд с маркированным списком"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    title_shape = slide.shapes.title
    title_shape.text = clean_markdown_text(title)
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    
    # Контент
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    
    # Ограничиваем количество пунктов
    display_bullets = bullets[:max_bullets]
    
    for i, bullet in enumerate(display_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(6)
        
        # Выделяем ключевые слова жирным
        if '**' in bullet or 'GigaChat' in bullet or 'Giga Web Insight' in bullet:
            p.font.bold = True
    
    # Если есть еще пункты, добавляем заметку
    if len(bullets) > max_bullets:
        p = tf.add_paragraph()
        p.text = f"... и еще {len(bullets) - max_bullets} пунктов"
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['light']
        p.font.italic = True
    
    return slide

def create_slide_with_table(prs, title, table_data):
    """Создает слайд с таблицей"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = clean_markdown_text(title)
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    
    # Создаем таблицу
    if table_data:
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows > 0 and cols > 0:
            # Ограничиваем размер таблицы для читаемости
            max_rows = min(rows, 8)
            max_cols = min(cols, 5)
            
            table = slide.shapes.add_table(
                max_rows, max_cols,
                Inches(0.5), Inches(1.2),
                Inches(9), Inches(4)
            ).table
            
            # Заполняем таблицу
            for i, row in enumerate(table_data[:max_rows]):
                for j, cell_text in enumerate(row[:max_cols]):
                    if i < max_rows and j < max_cols:
                        cell = table.cell(i, j)
                        cell.text = clean_markdown_text(cell_text)
                        cell.text_frame.paragraphs[0].font.size = Pt(11)
                        cell.text_frame.paragraphs[0].font.color.rgb = COLORS['text']
                        cell.text_frame.word_wrap = True
                        
                        # Заголовок таблицы (первая строка)
                        if i == 0:
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = COLORS['primary']
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    return slide

def create_title_slide(prs, title, subtitle=""):
    """Создает титульный слайд"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = clean_markdown_text(title)
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    if subtitle:
        subtitle_shape.text = clean_markdown_text(subtitle)
    else:
        subtitle_shape.text = "Интеграция GigaChat – бизнес‑кейс"
    
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['accent']
    
    return slide

def create_content_slide(prs, title, content_text):
    """Создает слайд с текстовым контентом"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = clean_markdown_text(title)
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = COLORS['primary']
    
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.word_wrap = True
    tf.text = clean_markdown_text(content_text)
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = COLORS['text']
    
    return slide

def should_combine_sections(section1, section2):
    """Определяет, стоит ли объединять два раздела"""
    # Объединяем короткие разделы
    total_content = len(section1.get("content", [])) + len(section2.get("content", []))
    total_subsections = len(section1.get("subsections", [])) + len(section2.get("subsections", []))
    
    # Объединяем если оба раздела короткие
    if total_content < 10 and total_subsections == 0:
        return True
    
    return False

def optimize_sections(sections):
    """Оптимизирует разделы, объединяя короткие"""
    optimized = []
    i = 0
    
    while i < len(sections):
        current = sections[i]
        
        # Пытаемся объединить с следующим разделом
        if i + 1 < len(sections) and should_combine_sections(current, sections[i + 1]):
            next_section = sections[i + 1]
            combined = {
                "title": f"{current['title']} / {next_section['title']}",
                "subsections": current.get("subsections", []) + next_section.get("subsections", []),
                "content": current.get("content", []) + next_section.get("content", [])
            }
            optimized.append(combined)
            i += 2
        else:
            optimized.append(current)
            i += 1
    
    return optimized

def convert_markdown_to_pptx(input_file, output_file=None):
    """Конвертирует Markdown файл в PowerPoint презентацию"""
    if output_file is None:
        # Генерируем имя выходного файла на основе входного
        import os
        base_name = os.path.splitext(os.path.basename(input_file))[0]
        output_file = f"{base_name}.pptx"
    
    # Читаем Markdown файл
    with open(input_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Парсим разделы
    sections = parse_markdown_sections(md_content)
    
    # Оптимизируем разделы
    sections = optimize_sections(sections)
    
    # Создаем титульный слайд
    if sections:
        main_title = "Giga Web Insight"
        create_title_slide(prs, main_title)
    
    # Обрабатываем разделы
    for section in sections:
        title = section["title"]
        content = section.get("content", [])
        subsections = section.get("subsections", [])
        
        # Обрабатываем раздел "Введение" отдельно
        if "Введение" in title:
            # Создаем слайд с миссией и продуктом
            intro_bullets = []
            for sub in subsections:
                if "Миссия" in sub['title']:
                    intro_bullets.append(f"🎯 {sub['title']}")
                    intro_bullets.extend(extract_bullets(sub['content'])[:2])
                elif "Продукт" in sub['title']:
                    intro_bullets.append(f"\n💡 {sub['title']}")
                    intro_bullets.extend(extract_bullets(sub['content'])[:5])
                elif "Рынок" in sub['title']:
                    intro_bullets.append(f"\n📊 {sub['title']}")
                    intro_bullets.extend(extract_bullets(sub['content'])[:2])
            
            if intro_bullets:
                create_slide_with_bullets(prs, "Введение", intro_bullets, max_bullets=10)
            elif content:
                bullets = extract_bullets(content)
                if bullets:
                    create_slide_with_bullets(prs, title, bullets)
            continue
        
        # Пропускаем титульный раздел
        if "Интеграция GigaChat" in title or not title:
            continue
        
        # Если есть подразделы, создаем отдельные слайды
        if subsections:
            # Для кейсов использования - группируем по 2 кейса на слайд
            if "Кейс" in title or "кейс" in title.lower():
                for i in range(0, len(subsections), 2):
                    if i + 1 < len(subsections):
                        # Два кейса на одном слайде
                        sub1 = subsections[i]
                        sub2 = subsections[i + 1]
                        combined_title = f"{title}"
                        combined_bullets = []
                        
                        # Кейс 1
                        combined_bullets.append(f"📌 {sub1['title']}")
                        bullets1 = extract_bullets(sub1['content'])
                        combined_bullets.extend(bullets1[:3])  # Первые 3 пункта
                        
                        # Разделитель
                        combined_bullets.append("")
                        
                        # Кейс 2
                        combined_bullets.append(f"📌 {sub2['title']}")
                        bullets2 = extract_bullets(sub2['content'])
                        combined_bullets.extend(bullets2[:3])  # Первые 3 пункта
                        
                        create_slide_with_bullets(prs, combined_title, combined_bullets, max_bullets=10)
                    else:
                        # Последний одиночный кейс
                        sub = subsections[i]
                        sub_title = f"{title}: {sub['title']}"
                        bullets = extract_bullets(sub['content'])
                        if bullets:
                            create_slide_with_bullets(prs, sub_title, bullets)
            else:
                # Для других разделов - по одному подразделу на слайд
                for sub in subsections:
                    sub_title = f"{title}: {sub['title']}"
                    bullets = extract_bullets(sub['content'])
                    if bullets:
                        create_slide_with_bullets(prs, sub_title, bullets)
        else:
            # Проверяем, есть ли таблица
            table_data = parse_table(content)
            if table_data:
                create_slide_with_table(prs, title, table_data)
            else:
                # Обычный слайд со списком
                bullets = extract_bullets(content)
                if bullets:
                    # Разбиваем на несколько слайдов если слишком много пунктов
                    max_per_slide = 6
                    for i in range(0, len(bullets), max_per_slide):
                        chunk = bullets[i:i+max_per_slide]
                        slide_title = title if i == 0 else f"{title} (продолжение)"
                        create_slide_with_bullets(prs, slide_title, chunk, max_per_slide)
                elif content:
                    # Текстовый слайд
                    content_text = '\n'.join(content[:5])  # Первые 5 строк
                    create_content_slide(prs, title, content_text)
    
    # Сохраняем презентацию
    prs.save(output_file)
    return output_file, len(prs.slides)

def main():
    """Основная функция для CLI использования"""
    import sys
    import os
    
    if len(sys.argv) > 1:
        input_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None
    else:
        input_file = 'PRESENTATION.md'
        output_file = 'PRESENTATION.pptx'
    
    if not os.path.exists(input_file):
        print(f"❌ Ошибка: файл {input_file} не найден")
        sys.exit(1)
    
    try:
        output_file, slide_count = convert_markdown_to_pptx(input_file, output_file)
        print(f"✅ Презентация создана: {output_file}")
        print(f"📊 Всего слайдов: {slide_count}")
        print(f"🎨 Использована цветовая схема: темно-синий (#003366)")
    except Exception as e:
        print(f"❌ Ошибка при создании презентации: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
